from pptx import Presentation
import comtypes.client
import os

def replace_placeholders(input_pptx_path, output_pptx_path, name, internship_name):
    """Replace placeholders in the PowerPoint presentation."""
    # Load the presentation
    presentation = Presentation(input_pptx_path)

    # Iterate through each slide and each shape
    for slide in presentation.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                # Replace placeholders in text
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = run.text.replace('<name>', name)
                        run.text = run.text.replace('<internship_name>', internship_name)

    # Save the modified presentation
    presentation.save(output_pptx_path)
    print(f"Modified PowerPoint saved as {output_pptx_path}")

def convert_pptx_to_pdf(pptx_path, pdf_path):
    """Convert a PowerPoint presentation to PDF."""
    # Initialize PowerPoint application
    powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
    powerpoint.Visible = 1

    # Open the presentation
    presentation = powerpoint.Presentations.Open(pptx_path)

    # Save as PDF
    presentation.SaveAs(pdf_path, 32)  # 32 = ppSaveAsPDF
    print(f"PDF saved as {pdf_path}")

    # Close the presentation and PowerPoint application
    presentation.Close()
    powerpoint.Quit()

if __name__ == "__main__":
    # User input for name and internship name
    user_name = input("Enter your name: ")
    internship_name = input("Enter your internship name: ")

    # Paths to input and output files
    input_pptx = r"C:\Users\91934\Downloads\template_letter.pptx"  # Update to your template path
    output_pptx = r"C:\Users\91934\Downloads\modified_letter.pptx"  # Temporary path for modified PPTX
    output_pdf = r"C:\new_projectr\letter_maker\output_letter.pdf"  # Output path for PDF

    # Replace placeholders in PowerPoint
    replace_placeholders(input_pptx, output_pptx, user_name, internship_name)

    # Convert the modified PowerPoint to PDF
    convert_pptx_to_pdf(output_pptx, output_pdf)

    # Optionally, delete the temporary PPTX file
    os.remove(output_pptx)
    print(f"Temporary PPTX file deleted.")
